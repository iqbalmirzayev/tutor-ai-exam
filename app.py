import streamlit as st
from PIL import Image
import cv2
import numpy as np
import shutil
# from streamlit.runtime.scriptrunner import add_script_run_context
import os
import threading
import json
from queue import Queue
import io
import zipfile
import uuid
import fitz  # PyMuPDF
from ultralytics import YOLO
from streamlit_drawable_canvas import st_canvas
from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.util import Inches as PptInches
from fpdf import FPDF
import tempfile

import asyncio
from aiogram import Bot
# from notifier import send_telegram_notification
# Telegram məlumatlarını bura daxil et
# Tokenləri Streamlit-in gizli secrets bölməsindən oxuyuruq
if "TELEGRAM_TOKEN" in st.secrets:
    TOKEN = st.secrets["TELEGRAM_TOKEN"]
    CHAT_ID = st.secrets["TELEGRAM_CHAT_ID"]
else:
    # Lokalda test edəndə xəta verməsin deyə (və ya bura öz tokenini müvəqqəti yaza bilərsən)
    TOKEN = "BOŞ"
    CHAT_ID = "BOŞ"

async def _async_send_notification(message):
    """Asinxron bildiriş göndərmə funksiyası."""
    bot = Bot(token=TOKEN)
    try:
        await bot.send_message(chat_id=CHAT_ID, text=message, parse_mode="Markdown")
    finally:
        # Sessiyanı bağlamaq vacibdir (yaxşı vərdiş!)
        await bot.session.close()

def send_telegram_notification(message):
    """Streamlit daxilində çağırmaq üçün təhlükəsiz sinxron körpü."""
    if TOKEN == "BOŞ" or TOKEN is None:
        return
    try:
        # Yeni bir hadisə döngəsi yaradırıq
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        loop.run_until_complete(_async_send_notification(message))
        loop.close()
    except Exception as e:
        print(f"Bildiriş xətası: {e}")
# --- 1. SƏHİFƏ TƏNZİMLƏMƏSİ ---
st.set_page_config(page_title="TutorAI", layout="wide")

# --- CSS (DİZAYN) ---
st.markdown("""
    <style>
    /* 1. Ümumi Arxa Fon və Şrift */
    .stApp {
        background-color: #0E1117;
    }
    
    /* 2. Şəkillərin Dizaynı */
    .stImage img { 
        border-radius: 12px; 
        box-shadow: 0 4px 6px rgba(0,0,0,0.3); 
        max-height: 300px !important; 
        object-fit: contain;
    }

    /* 3. Düymələrin Dizaynı (Daha modern, hover effekti ilə) */
    div.stButton > button { 
        width: 100%; 
        border-radius: 8px; 
        height: 40px; 
        font-weight: 600; 
        border: 1px solid #444; 
        transition: all 0.3s ease;
    }
    div.stButton > button:hover { 
        border-color: #00D4FF; 
        color: #00D4FF; 
        background-color: #262730;
    }
    
    /* 4. "Sil" düyməsi üçün xüsusi rəng (Qırmızımtıl) */
    div.stButton > button:active {
        transform: scale(0.98);
    }

    /* 5. İnput Qutuları (Səhifə nömrəsi yazılan yer) */
    div[data-testid="stNumberInput"] input { 
        text-align: center; 
        font-weight: bold; 
        border-radius: 8px;
    }

    /* 6. Sidebar (Yan Panel) Dizaynı */
    [data-testid="stSidebar"] {
        background-color: #161B22;
        border-right: 1px solid #30363D;
    }
    
    /* 7. Streamlit-in standart Header və Footer-ini gizlət */
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    header {visibility: hidden;}
    
    /* 8. Kart Dizaynı (Sual qutuları) */
    [data-testid="stVerticalBlock"] > [style*="flex-direction: column;"] > [data-testid="stVerticalBlock"] {
        /* Bu hissə konteynerlərə aiddir, border=True olanda işləyir */
    }
    </style>
""", unsafe_allow_html=True)

# --- AYARLAR ---
MODEL_PATH = "best.pt"
CANVAS_MAX_WIDTH = 800  
STROKE_COLOR = "#FF0000"
STROKE_WIDTH = 3

# --- MODELİ KEŞLƏ ---
@st.cache_resource
def load_model():
    return YOLO(MODEL_PATH)

def cleanup_old_sessions(base_dir="sessions", max_age_hours=24):
    import time
    if not os.path.exists(base_dir): return
    now = time.time()
    for folder in os.listdir(base_dir):
        folder_path = os.path.join(base_dir, folder)
        if os.path.getmtime(folder_path) < now - (max_age_hours * 3600):
            try:
                shutil.rmtree(folder_path)
            except: pass

if 'file_key' not in st.session_state:
    st.session_state.file_key = 0

if 'uploaded_pdf' not in st.session_state:
    st.session_state.uploaded_pdf = None

cleanup_old_sessions() # Proqram başlayanda işə düşür

def background_analyzer(user_dir, page_list):
    """Arxa planda verilmiş səhifələri analiz edib JSON kimi yadda saxlayır."""
    for p_idx in page_list:
        json_path = os.path.join(user_dir, f"results_{p_idx}.json")
        img_path = os.path.join(user_dir, f"page_{p_idx}.png")
        
        # Əgər bu səhifə artıq analiz olunubsa, keç
        if os.path.exists(json_path):
            continue
            
        if os.path.exists(img_path):
            img = cv2.imread(img_path)
            img_rgb = cv2.cvtColor(img, cv2.COLOR_BGR2RGB)
            
            # AI Analizi
            results = model.predict(img_rgb, conf=0.5, iou=0.45, verbose=False)[0]
            detected_boxes = results.boxes.data.tolist()
            
            raw_boxes = []
            for db in detected_boxes:
                if len(db) >= 4:
                    raw_boxes.append([int(x) for x in db[:4]])
            
            # Filtrləmə və Sıralama
            final_boxes = filter_overlapping_boxes(raw_boxes, iou_threshold=0.3)
            sorted_boxes = sort_boxes_column_wise(final_boxes, x_threshold=50)
            
            # Nəticəni JSON kimi diskə yazırıq
            with open(json_path, 'w') as f:
                json.dump(sorted_boxes, f)

try:
    model = load_model()
except Exception:
    st.error(f"❌ '{MODEL_PATH}' tapılmadı! Faylı qovluğa əlavə et.")
    st.stop()

# --- SESSİYA YADDAŞI ---
if 'ALL_QUESTIONS' not in st.session_state: st.session_state['ALL_QUESTIONS'] = {} 
if 'CURRENT_PAGE_IDX' not in st.session_state: st.session_state['CURRENT_PAGE_IDX'] = 0
if 'CANVAS_REFRESH_KEYS' not in st.session_state: st.session_state['CANVAS_REFRESH_KEYS'] = {}

# --- YARDIMÇI FUNKSİYALAR ---

def filter_overlapping_boxes(boxes, iou_threshold=0.3):
    if not boxes: return []
    boxes = np.array(boxes)
    if len(boxes) == 0: return []

    x1 = boxes[:, 0]
    y1 = boxes[:, 1]
    x2 = boxes[:, 2]
    y2 = boxes[:, 3]

    area = (x2 - x1) * (y2 - y1)
    idxs = np.argsort(area)
    pick = []

    while len(idxs) > 0:
        last = len(idxs) - 1
        i = idxs[last]
        pick.append(i) 
        xx1 = np.maximum(x1[i], x1[idxs[:last]])
        yy1 = np.maximum(y1[i], y1[idxs[:last]])
        xx2 = np.minimum(x2[i], x2[idxs[:last]])
        yy2 = np.minimum(y2[i], y2[idxs[:last]])
        w = np.maximum(0, xx2 - xx1)
        h = np.maximum(0, yy2 - yy1)
        intersection = w * h
        union = area[i] + area[idxs[:last]] - intersection
        iou = intersection / (union + 1e-6)
        idxs = np.delete(idxs, np.concatenate(([last], np.where(iou > iou_threshold)[0])))

    return boxes[pick].astype(int).tolist()

def get_page_image_from_disk(page_num):
    temp_dir = st.session_state.get('TEMP_DIR')
    if not temp_dir: return None
    image_path = os.path.join(temp_dir, f"page_{page_num}.png")
    if os.path.exists(image_path):
        # Faylı birbaşa oxuyub RGB-yə çeviririk
        img = cv2.imread(image_path)
        if img is None: return None
        return cv2.cvtColor(img, cv2.COLOR_BGR2RGB)
    return None

def sort_boxes_column_wise(boxes, x_threshold=50):
    if not boxes: return []
    boxes_sorted_x = sorted(boxes, key=lambda b: b[0])
    columns = []
    current_col = [boxes_sorted_x[0]]
    for i in range(1, len(boxes_sorted_x)):
        box = boxes_sorted_x[i]
        prev_box = boxes_sorted_x[i-1]
        if (box[0] - prev_box[0]) > x_threshold:
            columns.append(current_col)
            current_col = []
        current_col.append(box)
    columns.append(current_col)
    final_sorted = []
    for col in columns:
        col.sort(key=lambda b: b[1])
        final_sorted.extend(col)
    return final_sorted

def swap_questions(page_idx, idx1, idx2):
    boxes = st.session_state['ALL_QUESTIONS'][page_idx]
    boxes[idx1], boxes[idx2] = boxes[idx2], boxes[idx1]
    st.session_state['ALL_QUESTIONS'][page_idx] = boxes
    if page_idx not in st.session_state['CANVAS_REFRESH_KEYS']: st.session_state['CANVAS_REFRESH_KEYS'][page_idx] = 0
    st.session_state['CANVAS_REFRESH_KEYS'][page_idx] += 1

def delete_question(page_idx, idx):
    boxes = st.session_state['ALL_QUESTIONS'][page_idx]
    del boxes[idx]
    st.session_state['ALL_QUESTIONS'][page_idx] = boxes
    if page_idx not in st.session_state['CANVAS_REFRESH_KEYS']: st.session_state['CANVAS_REFRESH_KEYS'][page_idx] = 0
    st.session_state['CANVAS_REFRESH_KEYS'][page_idx] += 1

def process_image_for_export(img_rgb, invert=False):
    if invert:
        img_rgb = cv2.bitwise_not(img_rgb)
    return img_rgb

# --- UI BAŞLAYIR ---
st.title("🛠 TutorAI")
st.markdown("*Süni intellekt dəstəkli sual kəsmə və redaktə sistemi*", unsafe_allow_html=True)
st.divider()
if 'session_id' not in st.session_state:
    st.session_state['session_id'] = str(uuid.uuid4())

# 2. Qovluq strukturunu qururuq: sessions / <user_id>
base_dir = "sessions"
if not os.path.exists(base_dir):
    os.makedirs(base_dir)

user_dir = os.path.join(base_dir, st.session_state['session_id'])
if not os.path.exists(user_dir):
    os.makedirs(user_dir)

# Sessiyada qovluq yolunu yadda saxlayırıq ki, digər funksiyalar bilsin
st.session_state['TEMP_DIR'] = user_dir

# =========================================================================
# --- DÜZƏLİŞ EDİLƏN HİSSƏ (FILE UPLOADER & REFRESH LOGIC) ---
# =========================================================================

# 1. Əgər fayl yoxdursa -> Uploader-i göstər
if st.session_state.uploaded_pdf is None:
    uploaded_file = st.file_uploader(
        "Fayl yüklə (PDF tövsiyə olunur):", 
        type=["pdf", "jpg", "png"], 
        key=f"uploader_{st.session_state.file_key}" # Açar hər dəfə dəyişir
    )
    
    # Fayl seçilən kimi yaddaşa atıb səhifəni yeniləyirik
    if uploaded_file is not None:
        st.session_state.uploaded_pdf = uploaded_file
        st.rerun()
    
    # Fayl yoxdursa, aşağıdakı kodlar xəta verməməsi üçün dayandırırıq
    st.stop()

# 2. Fayl varsa -> Qutunu gizlət, əvəzinə Sil düyməsini göstər
else:
    # Faylı session_state-dən götürürük ki, kodun qalanı işləsin
    uploaded_file = st.session_state.uploaded_pdf

    col_info, col_del = st.columns([0.85, 0.15])
    with col_info:
        st.success(f"📂 Hazırda işlənən fayl: **{uploaded_file.name}**")
    with col_del:
        if st.button("❌ Sil", use_container_width=True):
            st.session_state.uploaded_pdf = None # Yaddaşı təmizlə
            st.session_state.file_key += 1       # Uploader-i sıfırla
            st.rerun()                           # Səhifəni yenilə

# =========================================================================
# --- BURADAN AŞAĞI HEÇ NƏ DƏYİŞMƏYİB (Orjinal Kod) ---
# =========================================================================

if uploaded_file:
    # Faylı unikal etmək üçün ad + ölçü + sessiya ID-sini birləşdiririk
    
    file_id = f"{uploaded_file.name}_{uploaded_file.size}"
    
    if st.session_state.get('LAST_FILE_ID') != file_id:
        send_telegram_notification(f"📢 *TutorAI istifadə edildi!*\n\n📄 Fayl: `{uploaded_file.name}`")
        # 1. KÖHNƏ DATA-NIN TƏMİZLƏNMƏSİ
        st.session_state['ALL_QUESTIONS'] = {}
        st.session_state['CURRENT_PAGE_IDX'] = 0
        st.session_state['LAST_FILE_ID'] = file_id
        
        if 'EXPORT_FILES' in st.session_state:
            del st.session_state['EXPORT_FILES']

        # 2. DİSKDƏKİ KÖHNƏ FAYLLARI SİLİRİK
        for f in os.listdir(user_dir):
            try:
                os.remove(os.path.join(user_dir, f))
            except: pass

        with st.spinner("📂 Yeni fayl hazırlanır..."):
            file_bytes = uploaded_file.getvalue()
            
            # PDF Emalı
            if uploaded_file.type == "application/pdf":
                doc = fitz.open(stream=file_bytes, filetype="pdf")
                total_p = len(doc) 
                st.session_state['TOTAL_PAGES'] = total_p
                
                # İlk 2 səhifəni dərhal emal et
                initial_pages = min(2, len(doc))
                for i in range(initial_pages):
                    page = doc[i]
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                    pix.save(os.path.join(user_dir, f"page_{i}.png"))
                
                def start_background_tasks(u_dir, f_bytes, total_pages_val):
                    initial_pages_inner = min(2, total_pages_val)
                    with fitz.open(stream=f_bytes, filetype="pdf") as d:
                        for i in range(initial_pages_inner, len(d)):
                            p = d[i]
                            p.get_pixmap(matrix=fitz.Matrix(2, 2)).save(os.path.join(u_dir, f"page_{i}.png"))
                    
                    # İlk 5 səhifəni analizə göndəririk
                    pages_to_analyze_inner = list(range(min(5, total_pages_val)))
                    background_analyzer(u_dir, pages_to_analyze_inner)

                threading.Thread(
                    target=start_background_tasks, 
                    args=(user_dir, file_bytes, total_p)
                ).start()
            
            # Şəkil emalı (PDF deyilsə)
            else:
                st.session_state['TOTAL_PAGES'] = 1
                nparr = np.frombuffer(file_bytes, np.uint8)
                img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
                cv2.imwrite(os.path.join(user_dir, "page_0.png"), img)
                threading.Thread(target=background_analyzer, args=(user_dir, [0])).start()
                
        st.rerun()

    total_pages = st.session_state.get('TOTAL_PAGES', 1)

    # --- NAVİQASİYA ---
    col_prev, col_nav, col_next = st.columns([1, 2, 1])
    
    if st.session_state['CURRENT_PAGE_IDX'] < 0: st.session_state['CURRENT_PAGE_IDX'] = 0
    elif st.session_state['CURRENT_PAGE_IDX'] >= total_pages: st.session_state['CURRENT_PAGE_IDX'] = total_pages - 1

    with col_prev:
        if st.button("⬅️ Öncəki", disabled=(st.session_state['CURRENT_PAGE_IDX'] == 0)):
            st.session_state['CURRENT_PAGE_IDX'] -= 1
            st.rerun()

    with col_nav:
        selected_page = st.number_input(
            "Səhifə", min_value=1, max_value=total_pages, 
            value=st.session_state['CURRENT_PAGE_IDX'] + 1,
            label_visibility="collapsed"
        )
        if selected_page - 1 != st.session_state['CURRENT_PAGE_IDX']:
            st.session_state['CURRENT_PAGE_IDX'] = selected_page - 1
            st.rerun()
        st.markdown(f"<div style='text-align: center; color: gray;'>Cəmi {total_pages} səhifə</div>", unsafe_allow_html=True)

    with col_next:
        if st.button("Növbəti ➡️", disabled=(st.session_state['CURRENT_PAGE_IDX'] == total_pages - 1)):
            st.session_state['CURRENT_PAGE_IDX'] += 1
            curr = st.session_state['CURRENT_PAGE_IDX']
            
            if (curr + 1) % 5 == 0:
                next_batch = list(range(curr + 1, min(curr + 6, total_pages)))
                threading.Thread(target=background_analyzer, args=(user_dir, next_batch)).start()
            
            st.rerun()

    # --- YAN PANEL ---
    st.sidebar.title("📊 Anbar")
    total_q_count = sum(len(boxes) for boxes in st.session_state['ALL_QUESTIONS'].values())
    active_pages = len([k for k, v in st.session_state['ALL_QUESTIONS'].items() if len(v) > 0])
    
    st.sidebar.markdown(f"""
    <div class="sidebar-stats">
        <h2 style='margin:0; color:#FF4B4B;'>{total_q_count}</h2>
        <span style='font-size:14px;'>Ümumi Sual Sayı</span>
        <hr style='margin:10px 0; border-color:#555;'>
        <span style='font-size:14px;'>Hazır Səhifələr: {active_pages}</span>
    </div>
    """, unsafe_allow_html=True)
    st.sidebar.divider()
    sequential_numbering = st.sidebar.toggle("🔢 Ardıcıl nömrələmə", value=True)

    if total_q_count > 0:
        st.sidebar.divider()
        st.sidebar.subheader("📤 Çıxarış Ayarları")
        dark_mode = st.sidebar.toggle("🌙 Dark Mode (Inverse)", value=False, help="Şəkillərin rəngini tərsinə çevir.")

        # --- FAYLLARI HAZIRLA DÜYMƏSİ ---
        if st.sidebar.button("⚙️ Sınağı Endir", type="primary"):
            progress_bar = st.sidebar.progress(0)
            status_text = st.sidebar.empty()
            
            with st.spinner("Fayllar hazırlanır..."):
                all_export_images = []
                
                total_to_process = len(range(total_pages))
                for idx, page_num in enumerate(range(total_pages)):
                    progress = (idx + 1) / total_to_process
                    progress_bar.progress(progress)
                    status_text.text(f"Emal olunur: Səhifə {page_num + 1}")

                    if page_num in st.session_state['ALL_QUESTIONS']:
                        boxes = st.session_state['ALL_QUESTIONS'][page_num]
                        if not boxes: continue
                        current_img = get_page_image_from_disk(page_num)
                        
                        if current_img is not None:
                            for box in boxes:
                                x1, y1, x2, y2 = map(int, box)
                                h, w, _ = current_img.shape
                                x1, y1 = max(0, x1), max(0, y1)
                                x2, y2 = min(w, x2), min(h, y2)
                                crop = current_img[y1:y2, x1:x2]
                                all_export_images.append(process_image_for_export(crop, dark_mode))

                status_text.text("📦 Fayllar paketlənir...")

                # 1. WORD
                doc = Document()
                doc.add_heading('Sınaq Sualları', 0)
                for img in all_export_images:
                    img_stream = io.BytesIO()
                    Image.fromarray(img).save(img_stream, format="PNG")
                    img_stream.seek(0)
                    doc.add_picture(img_stream, width=Inches(6))
                    doc.add_paragraph(" ") 
                doc_io = io.BytesIO()
                doc.save(doc_io)
                doc_io.seek(0)

                # 2. PPTX
                prs = Presentation()
                for img in all_export_images:
                    slide = prs.slides.add_slide(prs.slide_layouts[6]) 
                    img_stream = io.BytesIO()
                    Image.fromarray(img).save(img_stream, format="PNG")
                    img_stream.seek(0)
                    slide.shapes.add_picture(img_stream, PptInches(1), PptInches(1), width=PptInches(8))
                ppt_io = io.BytesIO()
                prs.save(ppt_io)
                ppt_io.seek(0)

                # 3. PDF
                pdf = FPDF()
                for img in all_export_images:
                    pdf.add_page()
                    with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
                        Image.fromarray(img).save(tmp.name)
                        tmp_path = tmp.name
                    pdf.image(tmp_path, x=10, y=10, w=190)
                    os.remove(tmp_path)
                pdf_bytes = pdf.output(dest='S').encode('latin-1') 

                # 4. ZIP
                zip_buffer = io.BytesIO()
                with zipfile.ZipFile(zip_buffer, "w") as zf:
                    for i, img in enumerate(all_export_images):
                        img_bgr = cv2.cvtColor(img, cv2.COLOR_RGB2BGR) 
                        success, encoded_img = cv2.imencode(".jpg", img_bgr)
                        if success:
                            zf.writestr(f"Sual_{i+1}.jpg", encoded_img.tobytes())
                
                st.session_state['EXPORT_FILES'] = {
                    "docx": doc_io,
                    "pptx": ppt_io,
                    "pdf": pdf_bytes,
                    "zip": zip_buffer.getvalue()
                }
            progress_bar.empty()
            status_text.empty()
            st.success("✅ Fayllar hazırdır! Aşağıdan yükləyə bilərsiniz.")

        # --- YÜKLƏMƏ DÜYMƏLƏRİ ---
        if 'EXPORT_FILES' in st.session_state:
            files = st.session_state['EXPORT_FILES']
            
            st.sidebar.download_button("📝 Word (.docx)", files["docx"], "sinaq.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
            st.sidebar.download_button("🖥️ PowerPoint (.pptx)", files["pptx"], "sinaq.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation")
            st.sidebar.download_button("🖨️ PDF (Sınaq)", files["pdf"], "sinaq.pdf", "application/pdf")
            st.sidebar.download_button("📦 ZIP (Şəkillər)", files["zip"], "sinaq_sekiller.zip", "application/zip")

    # --- İŞ MASASI ---
    st.divider()
    current_idx = st.session_state['CURRENT_PAGE_IDX']
    opencv_image = get_page_image_from_disk(current_idx)

    # Scaling
    orig_h, orig_w = opencv_image.shape[:2]
    if orig_w > CANVAS_MAX_WIDTH:
        scale_factor = CANVAS_MAX_WIDTH / orig_w
        new_w = CANVAS_MAX_WIDTH
        new_h = int(orig_h * scale_factor)
    else:
        scale_factor = 1.0
        new_w = orig_w
        new_h = orig_h
        
    resized_image = cv2.resize(opencv_image, (new_w, new_h))
    pil_image = Image.fromarray(resized_image)

    # --- AI ANALİZ / JSON OXUMA ---
    if current_idx not in st.session_state['ALL_QUESTIONS']:
        json_path = os.path.join(st.session_state['TEMP_DIR'], f"results_{current_idx}.json")
        
        if os.path.exists(json_path):
            try:
                with open(json_path, 'r') as f:
                    st.session_state['ALL_QUESTIONS'][current_idx] = json.load(f)
                
                if current_idx not in st.session_state['CANVAS_REFRESH_KEYS']: 
                    st.session_state['CANVAS_REFRESH_KEYS'][current_idx] = 0
                st.session_state['CANVAS_REFRESH_KEYS'][current_idx] += 1
                st.rerun()
            except Exception as e:
                pass

        with st.spinner(f"🔍 AI analiz edir... (Səhifə {current_idx + 1})"):
            results = model.predict(opencv_image, conf=0.5, iou=0.45, verbose=False)[0]
            detected_boxes = results.boxes.data.tolist()
            raw_boxes = []
            for db in detected_boxes:
                if len(db) >= 4:
                    raw_boxes.append([int(x) for x in db[:4]])
            
            final_filtered_boxes = filter_overlapping_boxes(raw_boxes, iou_threshold=0.3)
            sorted_boxes = sort_boxes_column_wise(final_filtered_boxes, x_threshold=50)
            
            st.session_state['ALL_QUESTIONS'][current_idx] = sorted_boxes
            
            with open(json_path, 'w') as f:
                json.dump(sorted_boxes, f)

            if current_idx not in st.session_state['CANVAS_REFRESH_KEYS']: 
                st.session_state['CANVAS_REFRESH_KEYS'][current_idx] = 0
            st.session_state['CANVAS_REFRESH_KEYS'][current_idx] += 1
            st.rerun()

    # --- CANVAS ---
    st.subheader(f"✏️ Səhifə {current_idx + 1}")
    
    alert_placeholder = st.empty()

    mode = st.radio("Rejim:", ("✋ Düzəliş", "➕ Yeni Sual"), horizontal=True, label_visibility="collapsed")
    drawing_mode = "transform" if mode == "✋ Düzəliş" else "rect"
    
    current_boxes = st.session_state['ALL_QUESTIONS'][current_idx]
    canvas_objects = []
    for box in current_boxes:
        x1, y1, x2, y2 = box
        canvas_objects.append({
            "type": "rect", "left": x1 * scale_factor, "top": y1 * scale_factor,
            "width": (x2 - x1) * scale_factor, "height": (y2 - y1) * scale_factor,
            "stroke": STROKE_COLOR, "strokeWidth": STROKE_WIDTH, "fill": "rgba(0,0,0,0)"
        })

    if current_idx not in st.session_state['CANVAS_REFRESH_KEYS']: st.session_state['CANVAS_REFRESH_KEYS'][current_idx] = 0
    refresh_count = st.session_state['CANVAS_REFRESH_KEYS'][current_idx]
    
    canvas_result = st_canvas(
        fill_color="rgba(0, 0, 0, 0)",
        stroke_width=STROKE_WIDTH, stroke_color=STROKE_COLOR,
        background_image=pil_image, update_streamlit=True,
        height=new_h, width=new_w, drawing_mode=drawing_mode,
        key=f"canvas_{st.session_state.get('LAST_FILE_ID', '')}_p{current_idx}_v{refresh_count}",
        initial_drawing={"version": "4.4.0", "objects": canvas_objects}
    )

    updated_boxes = []
    if canvas_result.json_data:
        for obj in canvas_result.json_data["objects"]:
            if obj["type"] == "rect":
                x = obj["left"] / scale_factor
                y = obj["top"] / scale_factor
                w = (obj["width"] * obj["scaleX"]) / scale_factor
                h = (obj["height"] * obj["scaleY"]) / scale_factor
                updated_boxes.append([int(x), int(y), int(x+w), int(y+h)])
    
    has_changes = False
    if len(updated_boxes) != len(current_boxes): has_changes = True
    elif updated_boxes != current_boxes: has_changes = True

    if has_changes:
        alert_placeholder.info("ℹ️ Diqqət: Dəyişiklik edilən zaman aşağıdakı 'Yenilə' düyməsini mütləq sıxın.")
        col_btn, _ = st.columns([1, 4])
        with col_btn:
            if st.button("🔄 Yenilə", type="primary"):
                st.session_state['ALL_QUESTIONS'][current_idx] = sort_boxes_column_wise(updated_boxes, x_threshold=50)
                st.session_state['CANVAS_REFRESH_KEYS'][current_idx] += 1
                st.rerun()

    st.write("---")

    final_boxes = st.session_state['ALL_QUESTIONS'].get(current_idx, [])

    start_num = 0
    if sequential_numbering:
        for p_idx in range(current_idx):
            start_num += len(st.session_state['ALL_QUESTIONS'].get(p_idx, []))

    if len(final_boxes) > 0:
        for i in range(0, len(final_boxes), 2):
            row_cols = st.columns(2)
            for j in range(2):
                if i + j < len(final_boxes):
                    with row_cols[j]:
                        idx = i + j
                        box = final_boxes[idx]
                        x1, y1, x2, y2 = map(int, box)
                        
                        x1, y1 = max(0, x1), max(0, y1)
                        x2, y2 = min(orig_w, x2), min(orig_h, y2)
                        crop = opencv_image[y1:y2, x1:x2]
                        
                        display_number = (start_num + idx + 1) if sequential_numbering else (idx + 1)
                        
                        with st.container(border=True):
                            c_ctrl, c_img = st.columns([1, 5])
                            
                            with c_ctrl:
                                st.markdown(f"<h3 style='text-align: center; color: #FF4B4B;'>{display_number}</h3>", unsafe_allow_html=True)
                                st.divider()
                                
                                if idx > 0:
                                    if st.button("⬆️", key=f"u_{current_idx}_{idx}"): 
                                        swap_questions(current_idx, idx, idx-1)
                                        st.rerun()
                                
                                if idx < len(final_boxes) - 1:
                                    if st.button("⬇️", key=f"d_{current_idx}_{idx}"): 
                                        swap_questions(current_idx, idx, idx+1)
                                        st.rerun()
                                
                                if st.button("🗑️", key=f"r_{current_idx}_{idx}"): 
                                    delete_question(current_idx, idx)
                                    st.rerun()
                            
                            with c_img: 
                                if crop.shape[0] > 0 and crop.shape[1] > 0:
                                    img_rgb = Image.fromarray(crop)
                                    buf = io.BytesIO()
                                    img_rgb.save(buf, format="JPEG", quality=85)
                                    st.image(buf.getvalue(), use_column_width=True)
                                else:
                                    st.error("⚠️ Şəkil ölçüsü xətalıdır.")
    else:
        st.info("Bu səhifədə sual yoxdur.")